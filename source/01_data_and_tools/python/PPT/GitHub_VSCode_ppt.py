from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

def add_footer(slide):
    left = Cm(1)
    top = Cm(17.2)
    width = Cm(23)
    height = Cm(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "GitHub with VSCode — Author: David Thébault"
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(191, 191, 191)
    p.alignment = PP_ALIGN.CENTER

def add_slide(title, contents, image_path=None, image_width_cm=None, image_height_cm=None, scale=1.0, notes=None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content_box = slide.placeholders[1]
    text_frame = content_box.text_frame
    text_frame.clear()

    for para in contents:
        p = text_frame.add_paragraph()
        p.text = para
        p.font.size = Pt(18)

    if image_path and image_width_cm and image_height_cm:
        display_width = Cm(image_width_cm * scale)
        display_height = Cm(image_height_cm * scale)
        left = Cm(14 - display_width.cm / 2)
        top = Cm(4)
        slide.shapes.add_picture(image_path, left, top, display_width, display_height)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    add_footer(slide)

# Slide 1 – Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Using GitHub with VSCode"
slide.placeholders[1].text = "Git Integration, Collaboration, and Best Practices"
add_footer(slide)

# Slides content
slides = [
    ("Why GitHub + VSCode?",
     ["Smooth integration with GitHub", "Powerful UI for commits, branches, and pull requests", "Extensions available for deeper Git control"],
     None, None, None, None, None
    ),

    ("Installing Git & VSCode",
     ["Download Git from https://git-scm.com", "Install VSCode from https://code.visualstudio.com", "Ensure Git is in your PATH"],
     None, None, None, None, "Include a screenshot of VSCode about page and terminal running 'git --version'."
    ),

    ("Open a Project Folder",
     ["Use File > Open Folder to load your project directory", "Recommended: Open the root folder of your repo"],
     None, None, None, None, "Add a screenshot of File > Open Folder in VSCode."
    ),

    ("Initialize a Git Repository",
     ["Click 'Initialize Repository' in the Source Control panel if no Git repo exists", "Python Example:\n```python\nimport os\nos.system('git init')\n```"],
     None, None, None, None, "Screenshot: Source Control panel with 'Initialize Repository' button."
    ),

    ("View Source Control",
     ["Click the Git icon on the left sidebar", "View changes, commits, branches in one place"],
     None, None, None, None, "Add screenshot of Source Control panel with staged/unstaged files."
    ),

    ("Stage & Commit Changes",
     ["Use '+' icon to stage files", "Write a commit message and click ✓ to commit",
      "Python Example:\n```python\nimport os\nos.system('git add .')\nos.system('git commit -m \"Initial commit\"')\n```"],
     "./png/stageFile.png", 10, 4, 0.75,
     "Include screenshot of commit message input and stage area."
    ),

    ("Connect to GitHub",
     ["Sign in with GitHub via VSCode's GitHub extension", "You may be prompted to authorize via browser"],
     None, None, None, None, "Screenshot of GitHub sign-in modal."
    ),

    ("Push to GitHub",
     ["Click 'Publish Branch' or use terminal:", "git remote add origin ...", "git push -u origin main",
      "Python Example:\n```python\nimport os\nos.system('git remote add origin https://github.com/yourusername/your-repo.git')\nos.system('git push -u origin main')\n```"],
     None, None, None, None, "Screenshot: publish branch button and terminal with push command."
    ),

    ("Clone a GitHub Repository",
     ["From GitHub: copy repo URL", "In VSCode: Ctrl+Shift+P > Git: Clone",
      "Python Example:\n```python\nimport os\nos.system('git clone https://github.com/yourusername/your-repo.git')\n```"],
     None, None, None, None, "Screenshot of command palette with 'Git: Clone' selected."
    ),

    ("Branching in VSCode",
     ["Use bottom-left status bar to switch/create branches", "GitLens or built-in features show branch tree",
      "Python Example:\n```python\nimport os\nos.system('git checkout -b new-feature')\n```"],
     None, None, None, None, "Screenshot of branch selection dropdown in bottom bar."
    ),

    ("Create a Pull Request",
     ["Click 'Create Pull Request' from Source Control panel (GitHub extension)", "Fill title and description"],
     None, None, None, None, "Screenshot of PR creation dialog in VSCode."
    ),

    ("Review a Pull Request",
     ["Fetch PRs from remote", "Open in editor and comment directly on lines", "Approve or request changes"],
     None, None, None, None, "Screenshot of a PR review inside VSCode."
    ),

    ("GitHub Panel in VSCode",
     ["GitHub tab lets you see issues, PRs, and discussions", "Great for collaboration from within the IDE"],
     None, None, None, None, "Show GitHub tab in the sidebar and active issues."
    ),

    ("Working with GitLens",
     ["GitLens extension enhances Git in VSCode", "Blame view, history, file annotations, etc."],
     None, None, None, None, "Show GitLens sidebar panel and blame view."
    ),

    ("Resolve Merge Conflicts",
     ["VSCode provides visual tools to accept changes", "Highlight conflict markers, click Accept Incoming/Current"],
     None, None, None, None, "Screenshot of a conflict with action buttons."
    ),

    ("Use the Terminal for Git",
     ["Access terminal inside VSCode", "Run any git command (add, commit, push, etc.)",
      "Python Example:\n```python\nimport os\nos.system('git status')\n```"],
     None, None, None, None, "Show integrated terminal with git commands in use."
    ),

    ("Settings and Extensions",
     ["Search for 'Git' or 'GitHub' in extensions panel", "Configure user name, email in settings.json",
      "Python Example:\n```python\nimport json\nsettings = {\n    \"git.user\": \"Your Name\",\n    \"git.email\": \"your.email@example.com\"\n}\nwith open('settings.json', 'w') as f:\n    json.dump(settings, f)\n```"],
     None, None, None, None, "Screenshot of extensions tab with GitHub extension installed."
    ),

    ("Common Issues & Fixes",
     ["Authentication issues → Re-authenticate GitHub", "Push rejected → Pull first", "Permission denied → SSH key or token needed"],
     None, None, None, None, "Show error message example and link to fix."
    ),

    ("Best Practices",
     ["Commit frequently with meaningful messages", "Create branches for features", "Pull before pushing", "Review PRs before merging"],
     None, None, None, None, None
    ),

    ("Resources",
     ["https://code.visualstudio.com/docs", "https://docs.github.com", "https://learngitbranching.js.org", "https://git-scm.com/book"],
     None, None, None, None, None
    ),
]

for title, content, image_path, width_cm, height_cm, scale, notes in slides:
    add_slide(title, content, image_path, width_cm, height_cm, scale, notes)

prs.save("./documents/github_vscode_training.pptx")
print("Presentation 'github_vscode_training.pptx' generated successfully.")
