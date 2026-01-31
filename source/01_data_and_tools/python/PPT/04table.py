import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Exemple de DataFrame
df = pd.DataFrame({
    "Year": ["2020", "2021", "2022", "2023"],
    "Revenue (€M)": [120, 150, 180, 210]
})

# Créer une présentation PowerPoint
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # layout avec seulement un titre
slide = prs.slides.add_slide(slide_layout)

# Ajouter un titre
slide.shapes.title.text = "Revenue Summary"

# Dimensions du tableau
rows, cols = df.shape[0] + 1, df.shape[1]  # +1 pour l'en-tête
left = Inches(2)
top = Inches(2)
width = Inches(6)
height = Inches(1.5)

# Ajouter le tableau
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Définir la largeur des colonnes
for col in range(cols):
    table.columns[col].width = Inches(3)

# Remplir l'en-tête
for j, col_name in enumerate(df.columns):
    cell = table.cell(0, j)
    cell.text = col_name
    para = cell.text_frame.paragraphs[0]
    para.font.bold = True
    para.font.size = Pt(16)
    para.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)

# Remplir les lignes de données
for i, row in df.iterrows():
    for j, value in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = str(value)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(16)
        para.alignment = PP_ALIGN.CENTER

# Enregistrer le fichier
prs.save("./Presentations/04table.pptx")
print("Slide created from DataFrame: 'slide_from_dataframe.pptx'")
