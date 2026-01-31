from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Dimensions in EMUs (English Metric Units)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Create a blank presentation and set slide size
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Add a blank slide
blank_slide_layout = prs.slide_layouts[6]  # blank layout
slide = prs.slides.add_slide(blank_slide_layout)

# Dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# --- Create red rectangle (1/3 of width) ---
red_width = int(slide_width * (1 / 3))
left = top = 0
height = slide_height

red_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, red_width, height
)
red_box.fill.solid()
red_box.fill.fore_color.rgb = RGBColor(200, 0, 0)  # Deep red
red_box.line.fill.background()  # No border

# --- Add title text centered in red area ---
title_shape = slide.shapes.add_textbox(left, Inches(2.5), red_width, Inches(2))
text_frame = title_shape.text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = "My Presentation Title"
run.font.size = Pt(36)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 255, 255)  # White text
p.alignment = 1  # Centered

# --- Save the presentation ---
prs.save("./Presentations/03redWhite.pptx")
print("Slide saved as 'red_white_title_slide.pptx'")
