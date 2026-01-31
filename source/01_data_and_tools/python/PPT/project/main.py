from pptx import Presentation
from slides.slide1_intro import add_slide1
from slides.slide2_plot import add_slide2
from slides.slide3_table import add_slide3

prs = Presentation()

add_slide1(prs)
add_slide2(prs)
add_slide3(prs)

prs.save("final_presentation.pptx")
print("Presentation saved as 'final_presentation.pptx'")
