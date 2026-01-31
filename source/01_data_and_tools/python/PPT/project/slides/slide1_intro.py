from pptx.util import Inches

def add_slide1(prs):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "UV"
    slide.placeholders[1].text = "The new pip?"