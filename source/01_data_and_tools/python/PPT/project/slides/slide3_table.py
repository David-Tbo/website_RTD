import pandas as pd
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide3(prs):
    df = pd.DataFrame({
        "Year": ["2020", "2021"],
        "Revenue": ["120", "150"]
    })

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Revenue Table"

    table = slide.shapes.add_table(df.shape[0]+1, df.shape[1],
                                   Inches(2), Inches(2), Inches(6), Inches(1.5)).table
    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx).text = col_name
    for i, row in df.iterrows():
        for j, val in enumerate(row):
            table.cell(i+1, j).text = str(val)
