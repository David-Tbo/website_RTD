import matplotlib.pyplot as plt
from pptx.util import Inches

def add_slide2(prs):
    plt.plot([2020, 2021, 2022], [100, 150, 200])
    plt.title("Sales over time")
    plt.savefig("plot.png")
    plt.close()

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Sales Graph"
    slide.shapes.add_picture("plot.png", Inches(1), Inches(2), height=Inches(4))
