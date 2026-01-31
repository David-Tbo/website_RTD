from pptx import Presentation

prs = Presentation()
slide_layout = prs.slide_layouts[0]  # Title slide
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Welcome"
subtitle.text = "This is the subtitle"

# Save the presentation
prs.save("./Presentations/00titleSlide.pptx")
