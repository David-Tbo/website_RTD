import os
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Create a simple DataFrame
data = {
    "Year": [2020, 2021, 2022, 2023],
    "Sales": [120, 150, 180, 210]
}
df = pd.DataFrame(data)

# Plot the DataFrame
plt.figure(figsize=(6, 4))
plt.plot(df["Year"], df["Sales"], marker='o')
plt.title("Annual Sales")
plt.xlabel("Year")
plt.ylabel("Sales")
plt.grid(True)

# Save the plot as an image
image_path = "sales_plot.png"
plt.savefig(image_path)
plt.close()

# Create a new presentation
prs = Presentation()

# Add a slide with a title + content layout
slide_layout = prs.slide_layouts[5]  # Title only layout
slide = prs.slides.add_slide(slide_layout)

# Set the title
title = slide.shapes.title
title.text = "Sales Over the Years"

# Insert the image (chart)
left = Inches(1)
top = Inches(2)
height = Inches(4.5)

slide.shapes.add_picture(image_path, left, top, height=height)

# Save the presentation
prs.save("./Presentations/02Chart.pptx")
print("Presentation saved as 'slide_with_chart.pptx'")
