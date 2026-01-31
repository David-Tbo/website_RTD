import pptx


print('hello from pptx')

# Create a new presentation
# ==========================================

prs = pptx.Presentation()


# Create the initial title / subtitle slide:
# ==========================================

# The layout 0 that we have chosen allows to have a title and a subtitle

slide1_register = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(slide1_register)

# (text) Add a title to the slide

title1 = slide1.shapes.title
title1.text = "Title of the slide N°1"

# Same command but all in one
# title1 = slide1.shapes.title.text = "Title of the slide N°1"

# (text) Add a subtitle

subtitle = slide1.placeholders[1]
subtitle.text = "Subtitle of the slide N°1"



# Create the bullet point slide:
# ==========================================

slide2_register = prs.slide_layouts[1]

# Add it to the presentation:
slide2 = prs.slides.add_slide(slide2_register)

# Add the title to this slide
title2 = slide2.shapes.title
title2.text = "Title of the slide N°2"

# Create a bullet point box
bullet_point_box = slide2.shapes

# level 1: Define the first bullet point + text of level 1
bullet_points_lvl1 = bullet_point_box.placeholders[1]
bullet_points_lvl1.text = "First bullet point text"
#bullet_points_lvl1.level = 0

# level 2: Define the first bullet point + text of level 2
bullet_points_lvl2 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl2.text = "Text of the bullet point of level 2"
bullet_points_lvl2.level = 1

# level 3: Define the first bullet point + text of level 2
bullet_points_lvl3 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl3.text = "Text of the bullet point of level 3"
bullet_points_lvl3.level = 2


# Save the presentation
prs.save("./Presentations/01bulletPoints.pptx")
