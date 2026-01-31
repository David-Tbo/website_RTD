# Fix the unpacking issue: split content and code blocks properly
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Footer utility function
def add_footer(slide):
    left = Inches(0.5)
    top = Inches(6.8)
    width = Inches(9)
    height = Inches(0.3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Git & GitHub presentation    Author: David Thébault"
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(191, 191, 191)  # Gray -25%
    p.alignment = PP_ALIGN.CENTER

# Utility function to add a slide with title and content
def add_slide(title, contents, code_blocks=None):
    layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    content_box = slide.placeholders[1]
    text_frame = content_box.text_frame
    text_frame.clear()

    for para in contents:
        p = text_frame.add_paragraph()
        p.text = para
        p.font.size = Pt(18)

    if code_blocks:
        for code in code_blocks:
            p = text_frame.add_paragraph()
            p.text = f">>> {code}"
            p.level = 0
            p.font.name = "Courier New"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 30, 30)
            p.space_before = Pt(6)
            p.space_after = Pt(6)

    add_footer(slide)

# Slide 1 — Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title_shape = slide.shapes.title
title_shape.text = "Git & GitHub"
title_paragraph = title_shape.text_frame.paragraphs[0]
title_paragraph.font.name = "Aptos Display"
title_paragraph.font.size = Pt(60)
subtitle_shape = slide.placeholders[1]
subtitle_shape.text = "Version Control and Collaboration for Developers"
add_footer(slide)

# Slides 2 to 21
slides = [
    ("What is Git?", ["Git is a distributed version control system.", "It allows tracking changes in source code during software development."], None),
    ("Why Use Git?", ["Version control", "Collaboration", "History and backups", "Branching and merging"], None),
    ("Installing Git", ["To install Git on your system:"], ["sudo apt install git", "brew install git", "git --version"]),
    ("Git Configuration", ["Set up your identity for commits:"], ["git config --global user.name \"Your Name\"", "git config --global user.email \"you@example.com\""]),
    ("Creating a Repository", ["Start versioning a project by initializing a repository:"], ["git init"]),
    ("Cloning a Repository", ["Copy a remote repository to your local machine:"], ["git clone https://github.com/user/repo.git"]),
    ("Basic Workflow", ["1. Modify files", "2. Stage changes", "3. Commit changes"], ["git add .", "git commit -m \"Your message\""]),
    ("Checking Status", ["See which files are modified/staged:"], ["git status"]),
    ("Viewing History", ["Inspect commit history:"], ["git log"]),
    ("Branches", ["Branches allow parallel development."], ["git branch", "git checkout -b new-feature"]),
    ("Merging", ["Merge branches together:"], ["git checkout main", "git merge new-feature"]),
    ("Undoing Changes", ["Revert changes or commits:"], ["git checkout -- file.py", "git reset HEAD file.py"]),
    ("What is GitHub?", ["GitHub is a platform for hosting and collaborating on Git repositories."], None),
    ("Creating a GitHub Repo", ["Use GitHub UI or CLI to create a repo.", "Then push local repo to GitHub:"], ["git remote add origin https://github.com/user/repo.git", "git push -u origin main"]),
    ("Pull Requests", ["Suggest changes by opening a pull request.", "Allows code review and discussion before merging."], None),
    ("Forking Repositories", ["Create your own copy of someone else's repository."], None),
    ("Issues and Discussions", ["Track bugs, suggest features, and discuss implementation using GitHub Issues and Discussions."], None),
    ("GitHub Actions", ["Automate workflows such as testing or deployment with GitHub Actions."], None),
    ("Best Practices", ["Commit often", "Write clear messages", "Use branches wisely", "Pull before you push"], None),
    ("Resources", ["https://git-scm.com/doc", "https://docs.github.com", "https://learngitbranching.js.org"], None),
]

for title, content, code in slides:
    add_slide(title, content, code_blocks=code)

# Save the presentation
prs.save("./Presentations/git_github_presentation.pptx")
print("Presentation 'git_github_presentation.pptx' generated successfully.")
